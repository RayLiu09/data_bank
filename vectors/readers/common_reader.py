import logging
import os.path
from datetime import datetime
from pathlib import Path

from .base_reader import BaseReader
from ..models.document import Document

logger = logging.getLogger(__name__)


class CommonReader(BaseReader):
    """
    A CommonReader class that can read the content for given files with extensions including .txt | .json | .md |.doc |.docx
    .ppt |.pptx and so on, except for the PDF files.
    """

    def __init__(self):
        super().__init__()
        self.name = "CommonReader"
        self.description = "A common reader for .txt, .docx, and .pptx"
        self.extensions = [".txt", ".docx", ".pptx"]

    def load(self, file_name: str, file_dir: str, **kwargs) -> list[Document]:
        """
        Load the content of the file

        @param file_name: The name of the file
        @param file_dir: The path of the file
        @param contents: The content of the file
        """
        documents = []
        if file_name is None:
            logger.error("File name is required")
            return []
        if not os.path.exists(file_name):
            logger.error(f"File {file_name} does not exist")
            return []

        ext = Path(file_name).suffix.lower()
        if ext not in self.extensions:
            logger.error(f"File extension {ext} is not supported")
            return []

        if ext == ".txt":
            with open(file_name, "r", encoding="utf-8") as file:
                content = file.read()
            doc = Document(name=file_name, ext=ext, content=content, timestamp=str(datetime.now().strftime("%Y-%m-%d "
                                                                                                           "%H:%M:%S")))
            documents.append(doc)

        if ext == ".docx":
            content = self._read_doc(file_name)
            if content is not None and len(content) > 0:
                doc = Document(name=file_name, ext=ext, content=content,
                               timestamp=str(datetime.now().strftime("%Y-%m-%d "
                                                                     "%H:%M:%S")))
                documents.append(doc)

        if ext == ".pptx":
            content = self._read_ppt(file_name)
            if content is not None and len(content) > 0:
                doc = Document(name=file_name, ext=ext, content=content,
                               timestamp=str(datetime.now().strftime("%Y-%m-%d "
                                                                     "%H:%M:%S")))
                documents.append(doc)

        logger.info(f"Loaded {len(documents)} documents from {file_name}")
        return documents

    def _read_doc(self, file_name: str) -> str:
        """
        Read the content of the .doc or .docx file
        """
        from docx import Document
        text = ""
        try:
            doc = Document(file_name)
            for para in doc.paragraphs:
                text += para.text + "\n"
            logger.info(f"Read {len(text)} characters from {file_name}")
        except Exception as e:
            logger.error(f"Failed to read {file_name}: {e}")
            raise ValueError("%s" % "Document format is not a valid docx format.")
        return text

    def _read_ppt(self, file_name: str) -> str:
        """
        Read the content of the .ppt or .pptx file
        """
        from pptx import Presentation
        text = ""
        prs = Presentation(file_name)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        logger.info(f"Read {len(text)} characters from {file_name}")
        return text
